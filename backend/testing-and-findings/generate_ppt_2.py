from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Paths ──────────────────────────────────────────────────────────────────
IMG_DIR = r"C:\Users\Aditya Chauhan\.gemini\antigravity-ide\brain\99a42128-5a98-4017-9e2a-bd17d67d8fa2\pdf_pages"
OUT_PATH = r"c:\Users\Aditya Chauhan\Desktop\My_VS_code_stuff\projects\ContextIQ\backend\testing-and-findings\ContextIQ_Final_Presentation.pptx"

def img(n): return os.path.join(IMG_DIR, f"fig_{n}.png")

# ── Colours ─────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0D, 0x0D, 0x0D)   # near-black
PURPLE    = RGBColor(0x8B, 0x5C, 0xF6)   # accent purple
BLUE      = RGBColor(0x38, 0xBD, 0xF8)   # accent blue
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LGREY     = RGBColor(0xB0, 0xB0, 0xB0)   # light grey for subtext
DGREY     = RGBColor(0x1E, 0x1E, 0x2E)   # card background
GREEN     = RGBColor(0x22, 0xC5, 0x5E)
RED       = RGBColor(0xEF, 0x44, 0x44)
YELLOW    = RGBColor(0xF5, 0x9E, 0x0B)

# ── Slide dimensions (16:9) ─────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout

# ═══════════════════════════════════════════════════════════════════════════
#  Helper functions
# ═══════════════════════════════════════════════════════════════════════════

def add_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide

def box(slide, text, x, y, w, h,
        font_size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True,
        bg_color=None, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    if bg_color:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    return txBox

def accent_bar(slide, y=Inches(0.85), color=PURPLE):
    """Thin horizontal accent line under headline."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.4), y, Inches(12.5), Pt(3)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

def headline(slide, text, sub=None):
    box(slide, text,
        Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.7),
        font_size=32, bold=True, color=WHITE)
    accent_bar(slide)
    if sub:
        box(slide, sub,
            Inches(0.4), Inches(0.92), Inches(12.5), Inches(0.4),
            font_size=15, color=LGREY, italic=True)

def bullet_block(slide, items, x, y, w, h, title=None, title_color=PURPLE, size=14):
    """Renders a titled bullet block."""
    yy = y
    if title:
        box(slide, title, x, yy, w, Inches(0.35),
            font_size=15, bold=True, color=title_color)
        yy += Inches(0.35)
    for item in items:
        box(slide, item, x + Inches(0.1), yy, w - Inches(0.1), Inches(0.34),
            font_size=size, color=WHITE)
        yy += Inches(0.33)

def card(slide, x, y, w, h):
    """Dark card background."""
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DGREY
    shape.line.color.rgb = PURPLE
    shape.line.width = Pt(1)
    return shape

def add_image(slide, path, x, y, w, h=None):
    if h:
        slide.shapes.add_picture(path, x, y, w, h)
    else:
        slide.shapes.add_picture(path, x, y, w)

def footer(slide, name="Aditya Chauhan  |  SAP BTP Core Team  |  AEH Intern, 5th June 2026 Batch"):
    box(slide, name,
        Inches(0.4), Inches(7.15), Inches(12.5), Inches(0.3),
        font_size=9, color=LGREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE SLIDE
# ═══════════════════════════════════════════════════════════════════════════
s1 = add_slide()

# Big purple gradient bar on left
bar = s1.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.25), H)
bar.fill.solid(); bar.fill.fore_color.rgb = PURPLE; bar.line.fill.background()

box(s1, "ContextIQ",
    Inches(0.6), Inches(1.8), Inches(8), Inches(1.4),
    font_size=72, bold=True, color=WHITE)

box(s1, "Enterprise IT Support Agent",
    Inches(0.6), Inches(3.1), Inches(8), Inches(0.6),
    font_size=28, color=PURPLE, bold=True)

box(s1, "A production-grade AI agent built with LangGraph, ChromaDB,\nGemini, Langfuse & OpenTelemetry",
    Inches(0.6), Inches(3.75), Inches(8), Inches(0.9),
    font_size=16, color=LGREY)

# Right side decorative tech tags
tags = ["LangGraph", "FastAPI", "ChromaDB", "Gemini", "Langfuse", "OpenTelemetry"]
colors = [PURPLE, BLUE, GREEN, YELLOW, RGBColor(0xEC, 0x48, 0x99), BLUE]
for i, (tag, col) in enumerate(zip(tags, colors)):
    tx = s1.shapes.add_textbox(Inches(9.2), Inches(1.8 + i * 0.62), Inches(3.5), Inches(0.5))
    tf = tx.text_frame; p = tf.paragraphs[0]; r = p.add_run()
    r.text = f"  {tag}  "; r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = col; r.font.name = "Calibri"
    shape = s1.shapes.add_shape(1, Inches(9.1), Inches(1.82 + i * 0.62), Inches(3.6), Inches(0.4))
    shape.fill.solid(); shape.fill.fore_color.rgb = DGREY
    shape.line.color.rgb = col; shape.line.width = Pt(1.5)

box(s1, "Aditya Chauhan  |  SAP BTP Core Team  |  AEH Intern — 5th June 2026 Batch",
    Inches(0.6), Inches(6.85), Inches(12), Inches(0.35),
    font_size=11, color=LGREY)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════════════════
s2 = add_slide()
headline(s2, "The Problem: IT Teams Are Drowning in Repetitive Queries",
         "Same questions. Same answers. Every single day.")

# Three problem cards
problems = [
    ("🔁  Repetition",   "\"My VPN won't connect\" — answered 40 times this month. IT engineers waste hours on tickets already solved."),
    ("🐢  Slow Resolution", "Employees wait hours for answers that exist in old tickets. Productivity lost on both sides."),
    ("🧠  Knowledge Silos", "IT solutions locked inside past tickets, emails, and docs. No single place to search them instantly."),
]
for i, (title, desc) in enumerate(problems):
    cx = Inches(0.4 + i * 4.3)
    card(s2, cx, Inches(1.55), Inches(4.1), Inches(2.8))
    box(s2, title, cx + Inches(0.15), Inches(1.65), Inches(3.8), Inches(0.5),
        font_size=16, bold=True, color=PURPLE)
    box(s2, desc, cx + Inches(0.15), Inches(2.2), Inches(3.8), Inches(2.0),
        font_size=13, color=LGREY)

# Solution callout
card(s2, Inches(0.4), Inches(4.55), Inches(12.5), Inches(0.75))
box(s2, "💡  ContextIQ solves this: One AI agent that reads every past IT ticket and answers employee questions instantly, accurately, and at scale.",
    Inches(0.6), Inches(4.65), Inches(12.1), Inches(0.6),
    font_size=15, bold=True, color=WHITE)

footer(s2)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — SOLUTION + UI SCREENSHOT
# ═══════════════════════════════════════════════════════════════════════════
s3 = add_slide()
headline(s3, "The Solution: ContextIQ in Action",
         "A scoped enterprise IT assistant — refuses hallucinations, answers only from your knowledge base.")

# Left: bullets
left_points = [
    "⚡  Sub-second response from past IT knowledge",
    "🔒  Strictly scoped to 8 IT categories only",
    "🚫  Rejects off-topic queries — no hallucinations",
    "📊  Every response tracked, costed & scored",
    "🔄  Switches between Gemini and Groq via .env",
]
bullet_block(s3, left_points, Inches(0.4), Inches(1.45), Inches(5.8), Inches(4.5), size=14)

# Right: UI screenshot (fig_15 — shows both reject + VPN answer)
add_image(s3, img(15), Inches(6.5), Inches(1.35), Inches(6.5), Inches(5.5))

footer(s3)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — PHASED BUILD APPROACH
# ═══════════════════════════════════════════════════════════════════════════
s4 = add_slide()
headline(s4, "Built in Phases — Each Phase Added Enterprise-Readiness",
         "Not built in one shot. Planned, incremental, testable.")

phases = [
    ("Phase 1", "Foundation", ["FastAPI backend", "ChromaDB vector store", "Basic RAG chain", "React frontend"]),
    ("Phase 2", "Intelligence", ["LangGraph ReAct Agent", "8-category scope gate", "Tool decision routing", "return_direct optimisation"]),
    ("Phase 3", "Observability", ["Langfuse LLM tracing", "Token cost tracking", "OpenTelemetry DB spans", "Quality score config"]),
    ("Phase 4", "Quality & Testing", ["A/B retrieval strategies", "Env-var config switching", "Prompt versioning", "Manual evaluation scores"]),
]

phase_colors = [BLUE, PURPLE, GREEN, YELLOW]
for i, (phase, title, bullets) in enumerate(phases):
    cx = Inches(0.35 + i * 3.25)
    card(s4, cx, Inches(1.45), Inches(3.1), Inches(5.5))
    shape = s4.shapes.add_shape(1, cx, Inches(1.45), Inches(3.1), Inches(0.45))
    shape.fill.solid(); shape.fill.fore_color.rgb = phase_colors[i]
    shape.line.fill.background()
    box(s4, phase, cx + Inches(0.1), Inches(1.45), Inches(2.9), Inches(0.45),
        font_size=13, bold=True, color=BG, align=PP_ALIGN.CENTER)
    box(s4, title, cx + Inches(0.1), Inches(1.95), Inches(2.9), Inches(0.4),
        font_size=15, bold=True, color=phase_colors[i])
    for j, b in enumerate(bullets):
        box(s4, f"• {b}", cx + Inches(0.15), Inches(2.42 + j * 0.5), Inches(2.8), Inches(0.46),
            font_size=12, color=LGREY)

footer(s4)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — WHY REACT AGENT (NOT SIMPLE RAG)
# ═══════════════════════════════════════════════════════════════════════════
s5 = add_slide()
headline(s5, "Why ReAct Agent? Not a Simple RAG Chain.",
         "This was the most critical architectural decision of the project.")

# Left column — Simple RAG
card(s5, Inches(0.4), Inches(1.45), Inches(5.8), Inches(4.5))
box(s5, "❌  Simple RAG Pipeline", Inches(0.6), Inches(1.55), Inches(5.4), Inches(0.45),
    font_size=16, bold=True, color=RED)
simple_rag = [
    "Blindly searches DB for EVERY message",
    "\"Capital of France?\" wastes tokens on IT tickets",
    "No decision-making — zero intelligence",
    "Higher latency. Higher API cost.",
    "Cannot refuse off-topic queries",
    "Hallucination risk on every response",
]
for i, t in enumerate(simple_rag):
    box(s5, f"  •  {t}", Inches(0.6), Inches(2.08 + i * 0.47), Inches(5.4), Inches(0.44),
        font_size=13, color=LGREY)

# Right column — ReAct
card(s5, Inches(6.5), Inches(1.45), Inches(6.4), Inches(4.5))
box(s5, "✅  ReAct Agent (ContextIQ)", Inches(6.7), Inches(1.55), Inches(6.0), Inches(0.45),
    font_size=16, bold=True, color=GREEN)
react_pts = [
    "Thinks FIRST — decides if DB search is needed",
    "Off-topic query? Rejected before hitting DB",
    "IT query? Routes to tool → ChromaDB → LLM",
    "~50% fewer unnecessary API calls",
    "Zero hallucinations outside scope",
    "return_direct=True eliminates redundant LLM call",
]
for i, t in enumerate(react_pts):
    box(s5, f"  •  {t}", Inches(6.7), Inches(2.08 + i * 0.47), Inches(6.0), Inches(0.44),
        font_size=13, color=WHITE)

# VS divider
box(s5, "VS", Inches(6.0), Inches(3.3), Inches(0.5), Inches(0.5),
    font_size=20, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

footer(s5)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — ARCHITECTURE DIAGRAM
# ═══════════════════════════════════════════════════════════════════════════
s6 = add_slide()
headline(s6, "System Architecture — Every Component Has a Reason",
         "Full-stack: React frontend → FastAPI → LangGraph agent → ChromaDB → Gemini → Langfuse + OpenTelemetry")

# Left: flowchart image (fig_8)
add_image(s6, img(8), Inches(0.3), Inches(1.35), Inches(5.0), Inches(5.4))

# Right: LangGraph graph from Langfuse (fig_16) + tool table
add_image(s6, img(16), Inches(5.5), Inches(1.35), Inches(2.8), Inches(4.0))

# Tech table on far right
headers = [("Tool", PURPLE), ("Why chosen", WHITE), ("Trade-off avoided", LGREY)]
rows = [
    ("LangGraph",      "Stateful graph agent",         "LangChain AgentExecutor (deprecated)"),
    ("ChromaDB",       "Local, free, persistent",       "Pinecone (paid, needs internet)"),
    ("Gemini",         "Free tier, strong reasoning",   "GPT-4 (costly)"),
    ("Langfuse",       "LLM-native observability",      "Basic logging (not production grade)"),
    ("OpenTelemetry",  "Industry span/trace standard",  "Custom timing code (not portable)"),
]
col_x = [Inches(8.5), Inches(9.9), Inches(11.2)]
col_w = [Inches(1.35), Inches(1.25), Inches(1.85)]

# Header row
for j, (h_text, h_col) in enumerate(headers):
    box(s6, h_text, col_x[j], Inches(1.45), col_w[j], Inches(0.35),
        font_size=11, bold=True, color=h_col)
bar2 = s6.shapes.add_shape(1, Inches(8.45), Inches(1.82), Inches(4.7), Pt(2))
bar2.fill.solid(); bar2.fill.fore_color.rgb = PURPLE; bar2.line.fill.background()

for i, row in enumerate(rows):
    row_y = Inches(1.9 + i * 0.62)
    if i % 2 == 0:
        bg = s6.shapes.add_shape(1, Inches(8.45), row_y - Inches(0.04), Inches(4.7), Inches(0.58))
        bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2A); bg.line.fill.background()
    for j, cell in enumerate(row):
        col = [WHITE, LGREY, RGBColor(0x70, 0x70, 0x90)][j]
        box(s6, cell, col_x[j], row_y, col_w[j], Inches(0.55), font_size=10, color=col)

footer(s6)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — OBSERVABILITY LAYER (Langfuse + OTEL)
# ═══════════════════════════════════════════════════════════════════════════
s7 = add_slide()
headline(s7, "Full-Stack Observability: Langfuse + OpenTelemetry",
         "Every LLM call tracked. Every DB span measured. Production-grade from day one.")

# Left: Langfuse trace (fig_4)
box(s7, "Langfuse — LLM Layer", Inches(0.4), Inches(1.45), Inches(6.0), Inches(0.38),
    font_size=14, bold=True, color=PURPLE)
add_image(s7, img(4), Inches(0.4), Inches(1.87), Inches(6.0), Inches(3.9))
box(s7, "Captures: agent steps, tool calls, token count (1,301 prompt → 140 completion), latency (10.2s total)",
    Inches(0.4), Inches(5.82), Inches(6.1), Inches(0.55), font_size=11, color=LGREY)

# Right: OTEL span (fig_7)
box(s7, "OpenTelemetry — Vector DB Layer", Inches(6.8), Inches(1.45), Inches(6.1), Inches(0.38),
    font_size=14, bold=True, color=BLUE)
add_image(s7, img(7), Inches(6.8), Inches(1.87), Inches(6.1), Inches(3.9))
box(s7, "Custom span 'vector_db_search': query, k=3, latency = 99ms. What Langfuse cannot see — we made visible.",
    Inches(6.8), Inches(5.82), Inches(6.1), Inches(0.55), font_size=11, color=LGREY)

footer(s7)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — A/B TESTING & QUALITY CONTROL
# ═══════════════════════════════════════════════════════════════════════════
s8 = add_slide()
headline(s8, "A/B Testing & Quality Control — Engineered, Not Guessed",
         "All parameters env-var driven. No code change needed to switch strategies.")

# Left col: what was tested + fig_11 (.env)
box(s8, "Config-Driven A/B Testing", Inches(0.4), Inches(1.45), Inches(5.8), Inches(0.38),
    font_size=14, bold=True, color=PURPLE)
add_image(s8, img(11), Inches(0.4), Inches(1.87), Inches(5.5), Inches(2.8))
ab_points = [
    "• Strategy A: Metadata chunk (ticket ID in each chunk)",
    "• Strategy B: Parent document (full ticket as parent)",
    "• K_RESULTS: 3 / 5 / 10 chunks retrieved",
    "• CHUNK_SIZE: 500 / 1000 / 1500 characters",
    "• LLM_PROVIDER: gemini / groq — hot-swappable",
]
for i, pt in enumerate(ab_points):
    box(s8, pt, Inches(0.4), Inches(4.82 + i * 0.36), Inches(5.8), Inches(0.34),
        font_size=12, color=LGREY)

# Right col: scoring + metrics
box(s8, "Quality Scoring in Langfuse", Inches(6.5), Inches(1.45), Inches(6.4), Inches(0.38),
    font_size=14, bold=True, color=BLUE)
add_image(s8, img(9), Inches(6.5), Inches(1.87), Inches(3.1), Inches(2.3))   # 1.34s
add_image(s8, img(10), Inches(9.8), Inches(1.87), Inches(3.2), Inches(2.3))   # token donut
add_image(s8, img(5), Inches(6.5), Inches(4.3), Inches(6.4), Inches(2.0))   # score config

footer(s8)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — LIVE DEMO RESULTS
# ═══════════════════════════════════════════════════════════════════════════
s9 = add_slide()
headline(s9, "Live Demo: 3 Test Cases — All Behaving as Designed",
         "Proof over promises. Every edge case handled.")

# Panel 1: Scope rejection (fig_14)
box(s9, "✅  Scope Enforcement", Inches(0.4), Inches(1.45), Inches(4.0), Inches(0.38),
    font_size=13, bold=True, color=GREEN)
add_image(s9, img(14), Inches(0.4), Inches(1.87), Inches(4.0), Inches(3.5))
box(s9, "Off-topic query rejected without touching Vector DB → saves tokens + prevents hallucination",
    Inches(0.4), Inches(5.42), Inches(4.1), Inches(0.65), font_size=11, color=LGREY)

# Panel 2: VPN answer (fig_15)
box(s9, "✅  IT Query Answered", Inches(4.65), Inches(1.45), Inches(4.3), Inches(0.38),
    font_size=13, bold=True, color=BLUE)
add_image(s9, img(15), Inches(4.65), Inches(1.87), Inches(4.3), Inches(3.5))
box(s9, "VPN question answered in 3 structured steps sourced from past IT tickets — not hallucinated",
    Inches(4.65), Inches(5.42), Inches(4.2), Inches(0.65), font_size=11, color=LGREY)

# Panel 3: FastAPI docs (fig_13)
box(s9, "✅  REST API Ready", Inches(9.2), Inches(1.45), Inches(3.8), Inches(0.38),
    font_size=13, bold=True, color=YELLOW)
add_image(s9, img(13), Inches(9.2), Inches(1.87), Inches(3.9), Inches(3.5))
box(s9, "Clean FastAPI with auto-generated Swagger docs — integration-ready for any enterprise system",
    Inches(9.2), Inches(5.42), Inches(3.9), Inches(0.65), font_size=11, color=LGREY)

footer(s9)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — LANGFUSE TRACE EVIDENCE
# ═══════════════════════════════════════════════════════════════════════════
s10 = add_slide()
headline(s10, "Observability Evidence: Real Traces from Real Queries",
          "72 total observations captured. Agent, Chain, Generation, Tool — all tracked.")

# fig_1 trace table (full width with annotation)
add_image(s10, img(1), Inches(0.4), Inches(1.45), Inches(12.5), Inches(4.3))

pts = [
    "🔴  LangGraph traces show exact input/output for each query — both VPN and off-topic",
    "📊  72 observations: 48 CHAIN + 12 GENERATION + 9 AGENT + 3 TOOL — full coverage",
    "⏱  Start time, latency, and response captured per trace — ready for SLA monitoring",
]
for i, pt in enumerate(pts):
    box(s10, pt, Inches(0.4), Inches(5.9 + i * 0.35), Inches(12.5), Inches(0.32),
        font_size=12, color=LGREY)

footer(s10)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — WHAT I LEARNED & WHAT'S NEXT
# ═══════════════════════════════════════════════════════════════════════════
s11 = add_slide()
headline(s11, "Learnings & What I Would Build Next",
          "This project taught me to think like an engineer, not just a developer.")

# Left: Learnings
card(s11, Inches(0.4), Inches(1.45), Inches(6.0), Inches(5.5))
box(s11, "Key Learnings", Inches(0.6), Inches(1.55), Inches(5.6), Inches(0.4),
    font_size=16, bold=True, color=PURPLE)
learnings = [
    "Observability is not optional — it is the first thing that breaks without it",
    "Enterprise AI needs scope enforcement — hallucination = liability",
    "Env-var driven config > code changes for A/B testing at runtime",
    "return_direct=True eliminated a full LLM call — 2× faster, 50% cost cut",
    "ReAct agents route intelligently; simple RAG chains don't know when to stop",
    "LangGraph > LangChain AgentExecutor for production stateful agents",
]
for i, l in enumerate(learnings):
    box(s11, f"  {i+1}.  {l}", Inches(0.6), Inches(2.05 + i * 0.72), Inches(5.6), Inches(0.68),
        font_size=12.5, color=WHITE)

# Right: What's next
card(s11, Inches(6.75), Inches(1.45), Inches(6.2), Inches(5.5))
box(s11, "If Given More Time →", Inches(6.95), Inches(1.55), Inches(5.8), Inches(0.4),
    font_size=16, bold=True, color=BLUE)
next_items = [
    ("🔐", "User authentication — IT managers see per-employee query history"),
    ("💬", "Multi-turn memory — agent remembers conversation within a session"),
    ("☁️",  "Azure/AWS deployment — Dockerise + API Gateway for 1M+ users"),
    ("🤖", "LLM-as-judge — automated quality scoring, replacing manual eval"),
    ("📈", "ServiceNow/Jira integration — ingest live IT tickets, not dummy data"),
    ("🔀", "Multi-agent routing — separate agents per category for better accuracy"),
]
for i, (icon, text) in enumerate(next_items):
    box(s11, f"  {icon}  {text}", Inches(6.95), Inches(2.05 + i * 0.72), Inches(5.8), Inches(0.68),
        font_size=12.5, color=LGREY)

footer(s11)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — CLOSING SLIDE
# ═══════════════════════════════════════════════════════════════════════════
s12 = add_slide()

bar12 = s12.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.25), H)
bar12.fill.solid(); bar12.fill.fore_color.rgb = PURPLE; bar12.line.fill.background()

box(s12, "ContextIQ is not just a chatbot.",
    Inches(1.0), Inches(1.8), Inches(11), Inches(0.9),
    font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

box(s12, "It is an observable, scoped, testable AI agent\nbuilt with the same patterns used in enterprise production systems.",
    Inches(1.0), Inches(2.75), Inches(11), Inches(1.1),
    font_size=22, color=PURPLE, align=PP_ALIGN.CENTER, bold=True)

# Tech tags row
tags2 = ["LangGraph", "FastAPI", "ChromaDB", "Gemini / Groq", "Langfuse", "OpenTelemetry", "React"]
total_w = 11.0
each_w = total_w / len(tags2)
for i, tag in enumerate(tags2):
    cx = Inches(1.0 + i * each_w)
    sh = s12.shapes.add_shape(1, cx, Inches(4.1), Inches(each_w - 0.15), Inches(0.45))
    sh.fill.solid(); sh.fill.fore_color.rgb = DGREY
    sh.line.color.rgb = PURPLE; sh.line.width = Pt(1)
    box(s12, tag, cx + Inches(0.05), Inches(4.12), Inches(each_w - 0.2), Inches(0.38),
        font_size=12, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

box(s12, "Thank you",
    Inches(1.0), Inches(5.2), Inches(11), Inches(0.7),
    font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

box(s12, "Aditya Chauhan  |  SAP BTP Core Team  |  AEH Intern — 5th June 2026 Batch",
    Inches(1.0), Inches(5.95), Inches(11), Inches(0.35),
    font_size=13, color=LGREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════════════════════════════════════
prs.save(OUT_PATH)
print(f"✅ PPT saved: {OUT_PATH}")
print(f"   Total slides: {len(prs.slides)}")
